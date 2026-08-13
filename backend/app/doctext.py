"""업로드된 문서에서 텍스트를 추출한다.

~/gitspace/mi-report의 pdftext.py/officetext.py 패턴(확장자별 추출기 + 실패 시
None/예외)을 참고했다. 그 프로젝트는 PyMuPDF(AGPL-3.0)·OCR·COM 폴백까지 갖춘
풀 파이프라인이지만, 여기서는 텍스트 기반 문서만 다룬다 — 스캔 PDF·OCR·xlsx는
이번 범위 밖이다(README '2차 확장 후보' 참고). pptx의 슬라이드 구분 표기
("--- slide n ---")는 mi-report의 officetext.py 출력 형식을 그대로 따랐다.
"""
from __future__ import annotations

import io

SUPPORTED_EXTENSIONS = (".txt", ".md", ".pdf", ".docx", ".pptx")


class UnsupportedDocumentError(ValueError):
    pass


def extract_text(filename: str, content: bytes) -> str:
    """파일명 확장자로 추출기를 고른다. 지원하지 않는 확장자면 예외를 낸다."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in (".txt", ".md"):
        return _extract_plain_text(content)
    if ext == ".pdf":
        return _extract_pdf(content)
    if ext == ".docx":
        return _extract_docx(content)
    if ext == ".pptx":
        return _extract_pptx(content)
    raise UnsupportedDocumentError(
        f"지원하지 않는 파일 형식입니다: {ext or filename} "
        f"(지원: {', '.join(SUPPORTED_EXTENSIONS)})"
    )


def _extract_plain_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p.strip() for p in pages if p.strip())


def _extract_docx(content: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(content))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for i, slide in enumerate(presentation.slides, start=1):
        lines = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if lines:
            slides.append(f"--- slide {i} ---\n" + "\n".join(lines))
    return "\n\n".join(slides)
